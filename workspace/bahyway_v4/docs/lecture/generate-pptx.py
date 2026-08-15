#!/usr/bin/env python3
"""
𒀭𒂗𒆤  BahyWay — EnkiDB Lecture PPTX Generator
================================================
Generates a fully-styled PowerPoint skeleton for the Iraqi Technology
University lecture on EnkiDB & Orbit-Oriented Ontology.

USAGE:
    pip install python-pptx
    python generate-pptx.py
    # → enkidb-lecture.pptx

REQUIRES:
    python-pptx >= 0.6.21
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx not installed.")
    print("Run:  pip install python-pptx")
    sys.exit(1)

# ── Sovereign Colours ─────────────────────────────────────────────────────────
BG      = RGBColor(0x0a, 0x0e, 0x1a)   # Deep navy background
GOLD    = RGBColor(0xd4, 0xa8, 0x43)   # Cuneiform gold
GEM     = RGBColor(0x00, 0xe5, 0xa0)   # GEM lane (inner ring)
TRIBE   = RGBColor(0x4d, 0xb8, 0xff)   # TRIBE lane
ACTIVE  = RGBColor(0xff, 0xd7, 0x00)   # ACTIVE lane
FUZZY   = RGBColor(0xff, 0x8c, 0x42)   # FUZZY lane
DEAD    = RGBColor(0x88, 0x88, 0x88)   # DEAD lane
WHITE   = RGBColor(0xe8, 0xe8, 0xe8)   # Body text
MUTED   = RGBColor(0x88, 0x88, 0x88)   # Subtitles
ARABIC  = RGBColor(0xe8, 0xd5, 0xa3)   # Arabic text

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # blank layout for full control


def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_bg(slide, color=BG):
    """Fill the slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def txbox(slide, text, left, top, width, height,
          size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          italic=False, font_name="Segoe UI"):
    """Add a text box and return the paragraph for further styling."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def heading(slide, text, top=Inches(0.3), size=36):
    txbox(slide, text,
          left=Inches(0.5), top=top,
          width=Inches(12.3), height=Inches(0.8),
          size=size, bold=True, color=GOLD)


def arabic_sub(slide, text, top=Inches(1.0)):
    txbox(slide, text,
          left=Inches(0.5), top=top,
          width=Inches(12.3), height=Inches(0.5),
          size=18, color=ARABIC, italic=True, align=PP_ALIGN.RIGHT)


def body(slide, text, top=Inches(1.6), size=20, color=WHITE, left=Inches(0.5), width=Inches(12.3), height=Inches(5.0)):
    txbox(slide, text, left=left, top=top, width=width, height=height,
          size=size, color=color)


def rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def divider(slide, top=Inches(1.15), color=GOLD):
    """Horizontal rule below heading."""
    ln = slide.shapes.add_shape(1, Inches(0.5), top, Inches(12.3), Pt(1))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()


# =============================================================================
# SLIDE 01 — TITLE
# =============================================================================
s = add_slide()
fill_bg(s)
# Cuneiform decoration
txbox(s, "𒀭𒂗𒆤", Inches(5.5), Inches(0.4), Inches(3), Inches(1.2),
      size=60, color=GOLD, align=PP_ALIGN.CENTER)
txbox(s, "EnkiDB", Inches(2), Inches(1.5), Inches(9), Inches(1.5),
      size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txbox(s, "Orbit-Oriented Ontology & the BahyWay Triple-O Philosophy",
      Inches(1.5), Inches(2.9), Inches(10), Inches(0.8),
      size=24, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "إنكي-قاعدة البيانات — فلسفة الكون المداري السيادي",
      Inches(1.5), Inches(3.6), Inches(10), Inches(0.6),
      size=20, color=ARABIC, align=PP_ALIGN.CENTER)
txbox(s, "Bahaa Fadam — BahyWay Sovereign Ecosystem",
      Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
      size=20, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txbox(s, "Iraqi Technology University, Baghdad · 2026",
      Inches(1.5), Inches(5.0), Inches(10), Inches(0.5),
      size=16, color=MUTED, align=PP_ALIGN.CENTER)
txbox(s, "v4.0.2 · 79 Crates · Pure Rust · 1,804 Tests · Zero External Databases",
      Inches(1.5), Inches(5.6), Inches(10), Inches(0.4),
      size=14, color=MUTED, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 02 — AGENDA
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "Agenda — 90 Minutes")
arabic_sub(s, "جدول الأعمال")
divider(s)
agenda = [
    ("I",   "The Problem",         "15 min", "Why all databases fail at scale"),
    ("II",  "Triple-O Philosophy", "20 min", "3 Axioms · KAKI · H(P) equation"),
    ("III", "The Engine",          "20 min", "12 Layers · BeeMDM Pipeline"),
    ("IV",  "Live Demo",           "15 min", "WASM Particle Visualizer"),
    ("V",   "Real Applications",   "10 min", "Najaf · Baghdad WPD · Nusku"),
    ("—",   "Q&A",                 "10 min", "Sovereignty Manifesto"),
]
rows = "\n".join(f"  {a[0]:<4}  {a[1]:<26}  {a[2]:<8}  {a[3]}" for a in agenda)
body(s, rows, top=Inches(1.7), size=19, color=WHITE)

# =============================================================================
# SLIDE 03 — SECTION BREAK: PART I
# =============================================================================
s = add_slide()
fill_bg(s, RGBColor(0x0a, 0x14, 0x2a))
txbox(s, "Part I",
      Inches(1), Inches(1.5), Inches(11), Inches(1.5),
      size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txbox(s, "The Problem",
      Inches(1), Inches(3.0), Inches(11), Inches(1.2),
      size=40, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "الجزء الأول — المشكلة",
      Inches(1), Inches(4.2), Inches(11), Inches(0.8),
      size=28, color=ARABIC, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 04 — THE CHAOS
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "Every Enterprise Has This Problem")
arabic_sub(s, "كل مؤسسة تعاني من هذه المشكلة")
divider(s)
body(s, (
    "The Spaghetti Reality:\n"
    "  • 12 databases, each with its own schema for 'customer'\n"
    "  • 3 ETL pipelines that disagree on what 'name' means\n"
    "  • Nobody knows which record is the truth\n"
    "  • DELETE destroys audit history\n"
    "  • UPDATE erases causality\n"
    "  • Soft-delete hacks everywhere\n\n"
    "\"A datum without an orbit is data without a home.\"\n"
    "بيانات بلا مدار هي بيانات بلا وطن"
), top=Inches(1.7), size=20)

# =============================================================================
# SLIDE 05 — SECTION BREAK: PART II
# =============================================================================
s = add_slide()
fill_bg(s, RGBColor(0x0a, 0x14, 0x2a))
txbox(s, "Part II", Inches(1), Inches(1.5), Inches(11), Inches(1.5),
      size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txbox(s, "Triple-O Philosophy", Inches(1), Inches(3.0), Inches(11), Inches(1.2),
      size=40, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "فلسفة الكون المداري المثلث", Inches(1), Inches(4.2), Inches(11), Inches(0.8),
      size=28, color=ARABIC, align=PP_ALIGN.CENTER)
txbox(s, "3 Axioms that replace all of relational theory",
      Inches(1), Inches(5.0), Inches(11), Inches(0.6),
      size=20, color=MUTED, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 06 — AXIOM 1
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "Axiom 1 — Every Datum is a Particle")
arabic_sub(s, "البديهية الأولى — كل بيانة هي جسيم")
divider(s)
body(s, (
    "A datum is NOT a row, document, or record.\n\n"
    "A datum is a PHYSICAL OBJECT with:\n"
    "  • Position in 7-dimensional quality space\n"
    "  • Velocity — rate of state change\n"
    "  • Density — neighbour count (SPH model)\n"
    "  • Orbital energy — quality score H(P)\n\n"
    "Physics metaphor: data obeys the same laws as matter.\n"
    "Every datum has a place it belongs,\n"
    "and we can measure how far it has drifted."
), top=Inches(1.7), size=20)

# =============================================================================
# SLIDE 07 — AXIOM 2: KAKI
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "Axiom 2 — Identity ≠ State (KAKI)")
arabic_sub(s, "البديهية الثانية — الهوية ليست الحالة")
divider(s)
body(s, (
    "KAKI Nucleus (κ) — 16 bytes — IMMUTABLE — born once, sealed forever\n\n"
    "  [ uuid_hash(4) | tribe_id(2) | type(1) | role(1) |\n"
    "    seq_counter(4) | epoch(2) | crc16(2) ]\n\n"
    "EAV Orbit — mutable · event-driven · append-only · time-stamped\n"
    "← All state lives here. NEVER in the nucleus.\n\n"
    "Three KAKI Types:\n"
    "  0x01  IdentityKaki  (KISHIB) — permanent sovereign particle\n"
    "  0x02  EventKaki     (ZIKRU)  — mutable state carrier\n"
    "  0x03  CrossTribeKaki(PARZU)  — inter-tribe connector\n\n"
    "\"The nucleus never changes. The orbit evolves.\"\n"
    "الجوهر لا يتغير. المدار يتطور."
), top=Inches(1.7), size=18)

# =============================================================================
# SLIDE 08 — AXIOM 3: H(P) EQUATION
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "Axiom 3 — Quality is Orbital Distance")
arabic_sub(s, "البديهية الثالثة — الجودة هي المسافة المدارية")
divider(s)
body(s, (
    "H(P) = 1 / (1 + √ Σ wᵢ(Pᵢ − Tᵢ)²)\n\n"
    "B11 = round(H(P) × 240)    →    [0 .. 240]\n\n"
    "7 Ibn Wahshiyya Dimensions:\n"
    "  D1 Accuracy     w=0.30    D5 Uniqueness   w=0.10\n"
    "  D2 Completeness w=0.20    D6 Timeliness   w=0.05\n"
    "  D3 Consistency  w=0.15    D7 Integrity    w=0.05\n"
    "  D4 Validity     w=0.15\n\n"
    "Quality Lanes:\n"
    "  GEM    B11 ≥ 200    Inner ring  (35.4% target)\n"
    "  TRIBE  B11 140-199  Mid ring\n"
    "  ACTIVE B11 100-139  Mid ring\n"
    "  FUZZY  B11  60- 99  Outer ring  (monitored)\n"
    "  DEAD   B11   < 60   Rule 7 sink\n\n"
    "Why 240 not 255? ← Plimpton 322 (Babylonian base-60)"
), top=Inches(1.7), size=17)

# =============================================================================
# SLIDE 09 — THE 7 TRIBAL LAWS
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "The 7 Tribal Laws")
arabic_sub(s, "القوانين القبلية السبعة")
divider(s)
body(s, (
    "Priority (ascending): P(L3) < P(L7) < P(L1) < P(L2) < P(L5) < P(L4) < P(L6)\n\n"
    "L3  Scope Law      — Domain boundary enforcement\n"
    "L7  Overflow Law   — Density ≥ 11 neighbours → DeadArchive (Rule 7)\n"
    "L1  Identity Law   — KAKI nucleus locked forever\n"
    "L2  Quality Law    — Promotes particles that improve H(P)\n"
    "L5  Transition Law — Stabilises FUZZY, prevents oscillation\n"
    "L4  Orbit Law      — Governs GEM ascent via centroid\n"
    "L6  Sovereignty Law — SUPREME: GEM particles protected\n\n"
    "Sovereign constants:\n"
    "  TAU_R7 = 11  (Rule 7 trigger)\n"
    "  RESONANCE_RADIUS = 0.15  (SPH kernel window)\n"
    "  GEM_RATE_TARGET = 35.4%"
), top=Inches(1.7), size=19)

# =============================================================================
# SLIDE 10 — SECTION BREAK: PART III
# =============================================================================
s = add_slide()
fill_bg(s, RGBColor(0x0a, 0x14, 0x2a))
txbox(s, "Part III", Inches(1), Inches(1.5), Inches(11), Inches(1.5),
      size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txbox(s, "The Engine", Inches(1), Inches(3.0), Inches(11), Inches(1.2),
      size=40, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "79 crates · 12 layers · Pure Rust · Zero external databases",
      Inches(1), Inches(4.2), Inches(11), Inches(0.6),
      size=22, color=MUTED, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 11 — 12-LAYER ARCHITECTURE
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "12-Layer Architecture")
arabic_sub(s, "البنية المعمارية ذات الـ ١٢ طبقة")
divider(s)
body(s, (
    "L12  bahyway-web           WASM · Canvas 2D · Browser target\n"
    "L11  dubsar-ide · dubsar-visualizer\n"
    "L10  eridu-runtime · scheduler · supervisor\n"
    "L9.5 enkidullm-* · zikru-embed      Sovereign LLM\n"
    "L9.1 kupru · akkvalue · istar        Crypto / Access Control\n"
    "L9   aaol · heptascript · akkadi     Languages\n"
    "L8   bahyway-fabric · bahyway-dqm    ◄ Enterprise Data Fabric\n"
    "L7   template-engine · damadmbok-dictionary\n"
    "L6   idu-prober · idu-batching\n"
    "L5   story · fuzzy · score · najaf · wpd · dmw (18 engines)\n"
    "L4.5 vgca-engine · tribe-orbit-engine · ammas-engine  Physics\n"
    "L4   enkidb-engine · enkidb-query\n"
    "L3   enkidb-indexes (7 types) · enkidb-dictionary\n"
    "L2   block · journal · storage · snapshot · quantdb\n"
    "L1   enkidb-kaki · enkidb-vector-id\n"
    "L0   bahyway-core · bahyway-crc · bahyway-algebra"
), top=Inches(1.7), size=15)

# =============================================================================
# SLIDE 12 — BEEMDM PIPELINE
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "BeeMDM — Sovereign ETL Pipeline")
arabic_sub(s, "بي-إم-دي-إم — خط أنابيب البيانات السيادي")
divider(s)
body(s, (
    "4 Lanes × 8 Sovereign Stations:\n\n"
    "S0  adad-gate              VaultGate — sole ingestion point · KAKI minter\n"
    "S1  musaru-security        Auth · Authorization · Threat detection\n"
    "S2  vgca-validation        VGCA 7D FSV + 6D BFV geometry validation\n"
    "S3  data-structure-station Structural conformance\n"
    "S4  data-cleansing-station VGCA transformation + DQM 6-dimension scoring\n"
    "S5  bahyway-dqm            DAMA-DMBOK 6 dimensions:\n"
    "                           Completeness · Validity · Accuracy\n"
    "                           Consistency · Uniqueness · Timeliness\n"
    "S6  idu-prober             Cross-tribe identity resolution\n"
    "S7  data-steward-station   Governance + approval workflow\n"
    "S8  permanent-storage      Golden Records · Immutable · CRC-16 sealed ✓"
), top=Inches(1.7), size=17)

# =============================================================================
# SLIDE 13 — STORY ENGINE
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "StoryEngine — Time-Travel Queries")
arabic_sub(s, "محرك القصة — استعلامات السفر عبر الزمن")
divider(s)
body(s, (
    "CQRS Architecture:\n"
    "  Write side: Append-only Journal (Event-KAKIs)\n"
    "  Read side:  Event projection (state always projected, never stored)\n\n"
    "  project_at(particle, epoch) algorithm:\n"
    "    1. Find snapshot ≤ target epoch\n"
    "    2. Load snapshot state\n"
    "    3. Apply delta Event-KAKIs since snapshot\n"
    "    4. Return complete state at epoch T\n\n"
    "No UPDATE. No DELETE. Ever.\n"
    "Every state change adds a new Event-KAKI.\n"
    "The old state is still there. Always.\n\n"
    "كل تغيير يضيف ختماً جديداً. الحالة القديمة لا تُمحى أبداً."
), top=Inches(1.7), size=20)

# =============================================================================
# SLIDE 14 — ORBITAL TRUST PROBE
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "orbital-trust-probe — Self-Regulating Quality")
arabic_sub(s, "مسبار الثقة المداري — الجودة ذاتية التنظيم")
divider(s)
body(s, (
    "Problem: At billion-particle scale, a particle may drift even\n"
    "when fuzzy rules haven't changed. Data quality or trust failure?\n\n"
    "4-Step Causal Attribution:\n"
    "  Step 1  FuzzyRules fingerprint — did rules change?\n"
    "  Step 2  StoryEngine EAV delta  — legitimate state evolution?\n"
    "  Step 3  ScoreEngine field      — neighbour density shift?\n"
    "  Step 4  Unexplained residual   → trust penalty applied\n\n"
    "Closed Feedback Loop:\n"
    "  OrbitalDeviation\n"
    "  → D9 penalty → penalised D6 (SourceTrust)\n"
    "  → lower B11 → organic ring correction\n"
    "  → self-terminating cascade ✓\n\n"
    "22 tests · OrbitalDeviationJournal (CRC-16 sealed, append-only)"
), top=Inches(1.7), size=19)

# =============================================================================
# SLIDE 15 — SECTION BREAK: LIVE DEMO
# =============================================================================
s = add_slide()
fill_bg(s, RGBColor(0x0a, 0x14, 0x2a))
txbox(s, "Part IV", Inches(1), Inches(1.5), Inches(11), Inches(1.5),
      size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txbox(s, "Live Demo", Inches(1), Inches(3.0), Inches(11), Inches(1.2),
      size=40, color=GEM, align=PP_ALIGN.CENTER)
txbox(s, "Rust → WebAssembly → Browser · bahyway.com",
      Inches(1), Inches(4.2), Inches(11), Inches(0.6),
      size=22, color=MUTED, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 16 — DEMO DESCRIPTION
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "BeeMDM Live Demo — What You Will See")
arabic_sub(s, "ماذا ستشاهد في العرض الحي")
divider(s)
body(s, (
    "Panel 1 — Orbital Rings\n"
    "  100 particles orbiting in real time.\n"
    "  Color = quality lane (GEM/TRIBE/ACTIVE/FUZZY/DEAD).\n"
    "  Watch particles migrate as quality scores change.\n\n"
    "Panel 2 — BeeMDM Pipeline\n"
    "  Live ETL processing — each station lights up as data flows.\n"
    "  DQM 6-dimension heatmap per record.\n\n"
    "Panel 3 — Story Engine\n"
    "  Time-travel slider. Drag to any epoch.\n"
    "  CQRS projection computed live in the browser.\n\n"
    "Tech Stack:\n"
    "  Rust compiled to WASM via wasm-bindgen\n"
    "  Canvas 2D rendering via web-sys bindings\n"
    "  Zero server calls · Runs entirely in your browser\n"
    "  يعمل بالكامل في متصفحك — لا خادم مطلوب"
), top=Inches(1.7), size=19)

# =============================================================================
# SLIDE 17 — SECTION BREAK: APPLICATIONS
# =============================================================================
s = add_slide()
fill_bg(s, RGBColor(0x0a, 0x14, 0x2a))
txbox(s, "Part V", Inches(1), Inches(1.5), Inches(11), Inches(1.5),
      size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txbox(s, "Real Applications", Inches(1), Inches(3.0), Inches(11), Inches(1.2),
      size=40, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "التطبيقات الحقيقية", Inches(1), Inches(4.2), Inches(11), Inches(0.8),
      size=28, color=ARABIC, align=PP_ALIGN.CENTER)

# =============================================================================
# SLIDE 18 — NAJAF ENGINE
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "NajafEngine — Cemetery Navigation")
arabic_sub(s, "محرك النجف — نظام ملاحة المقابر الكبرى")
divider(s)
body(s, (
    "Wadi Al-Salam, Najaf — the world's largest cemetery.\n"
    "5–6 million graves. Families cannot find ancestors.\n"
    "No unified registry. No spatial index.\n\n"
    "The EnkiDB Solution:\n"
    "  • Each grave = sovereign KAKI particle (IdentityKaki 0x01)\n"
    "  • Spatial position in H3 hexagonal index (D5 geo_precision)\n"
    "  • Family relations via CrossTribeKaki bridges\n"
    "  • Historical records via StoryEngine time-travel\n"
    "  • orbital-trust-probe detects archival data drift\n\n"
    "crate: najaf-engine\n"
    "binary: bin/najaf-ingest  ← sovereign ingestion pipeline"
), top=Inches(1.7), size=20)

# =============================================================================
# SLIDE 19 — WPD + NUSKU + AMMAS
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "WPD · DMW · Nusku · AMMAS")
arabic_sub(s, "أنابيب المياه · المشورة · الحرارة · الجسيمات الحركية")
divider(s)
body(s, (
    "WPDEngine — Baghdad Water Pipeline\n"
    "  Each pipe segment = KAKI particle\n"
    "  Pressure readings = Event-KAKIs\n"
    "  orbital-trust-probe detects anomalous pressure before leaks form\n\n"
    "DMWEngine — SQL Advisor\n"
    "  Translates legacy SQL queries → KAKI-based retrieval\n"
    "  Bridges relational systems to sovereign data architecture\n\n"
    "NuskuEngine — Thermal Screening\n"
    "  (Named after the Sumerian god of fire and light)\n"
    "  IR camera → particle orbit analysis · 50ms SLA pipeline budget\n\n"
    "AMMASEngine — Kinetic Particle System\n"
    "  (Named after the Akkadian war god)\n"
    "  SPH density + orbital velocity · models large-scale migrations"
), top=Inches(1.7), size=18)

# =============================================================================
# SLIDE 20 — COMPARISON TABLE
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "EnkiDB vs. Traditional Databases")
arabic_sub(s, "إنكي-قاعدة البيانات مقابل قواعد البيانات التقليدية")
divider(s)
body(s, (
    "Capability          PostgreSQL/MySQL              EnkiDB\n"
    "─────────────────────────────────────────────────────────\n"
    "Identity            AUTO_INCREMENT/UUID           KAKI 16-byte sovereign seal\n"
    "Data Quality        Constraints (syntax only)     H(P) equation, 7D orbital score\n"
    "Deletion            DELETE (destroys history)      Forbidden — INSERT supersedes\n"
    "Update              UPDATE (erases causality)      New Event-KAKI appended\n"
    "Time Travel         Audit log add-on               Native CQRS project_at(epoch)\n"
    "Trust               Assumed                        Computed — orbital probe\n"
    "Dependencies        External runtime required      Pure Rust — zero external\n"
    "Arabic Support      Collation hacks                VGCA-Σ 7D native geometry"
), top=Inches(1.7), size=17)

# =============================================================================
# SLIDE 21 — BY THE NUMBERS
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "BahyWay.Ecosystem v4.0 — By the Numbers")
arabic_sub(s, "النظام الإيكولوجي بالأرقام")
divider(s)
body(s, (
    "Sovereign Crates         79\n"
    "Passing Tests         1,804\n"
    "Lines of Rust        74,022\n"
    "Architecture Layers      12\n"
    "External Databases        0\n"
    "unsafe {} blocks (physics/intelligence)    0\n"
    "Bytes per KAKI           16\n"
    "KAKIs per tribe per epoch   4,294,967,294\n"
    "Forbidden Operations     17\n"
    "Eternal Constants         6  (QUALITY_DIVISOR=240, TAU_R7=11, …)"
), top=Inches(1.7), size=22)

# =============================================================================
# SLIDE 22 — SOVEREIGNTY MANIFESTO
# =============================================================================
s = add_slide()
fill_bg(s)
heading(s, "The Sovereignty Manifesto")
arabic_sub(s, "البيان السيادي")
divider(s)
body(s, (
    "Your data is sovereign territory.\n"
    "It obeys your laws, not Amazon's.\n\n"
    "Every datum deserves an immutable identity.\n"
    "KAKI is its permanent passport.\n\n"
    "Quality is not a checkbox.\n"
    "It is a mathematical distance in orbital space.\n\n"
    "History is sacred.\n"
    "No operation may erase what was.\n\n"
    "Memory safety is sovereignty.\n"
    "Rust enforces what laws cannot.\n\n"
    "Iraqi mathematics built this.\n"
    "Plimpton 322 · Ibn Wahshiyya · Enki · Gilgamesh.\n\n"
    "بياناتك أرض سيادية. هويتها أبدية. جودتها رياضية. تاريخها مقدس."
), top=Inches(1.7), size=19)

# =============================================================================
# SLIDE 23 — Q&A
# =============================================================================
s = add_slide()
fill_bg(s)
txbox(s, "𒀭𒂗𒆤", Inches(5.5), Inches(0.5), Inches(3), Inches(1.2),
      size=60, color=GOLD, align=PP_ALIGN.CENTER)
txbox(s, "Questions?",
      Inches(2), Inches(1.7), Inches(9), Inches(1.5),
      size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
txbox(s, "أسئلة؟",
      Inches(2), Inches(3.2), Inches(9), Inches(0.8),
      size=36, color=ARABIC, align=PP_ALIGN.CENTER)
txbox(s, "github.com/bahyway/enkidb  ·  bahyway.com  ·  heptascript.com",
      Inches(2), Inches(4.2), Inches(9), Inches(0.5),
      size=18, color=TRIBE, align=PP_ALIGN.CENTER)
txbox(s, "Bahaa Fadam  ·  bahaa.fadam@gmail.com",
      Inches(2), Inches(4.8), Inches(9), Inches(0.5),
      size=18, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "EnkiDB v4.0.2 · 79 Crates · Pure Rust · 1,804 Tests",
      Inches(2), Inches(5.4), Inches(9), Inches(0.4),
      size=14, color=MUTED, align=PP_ALIGN.CENTER)

# =============================================================================
# SAVE
# =============================================================================
out = "enkidb-lecture.pptx"
prs.save(out)
print(f"✓  Saved: {out}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Open with Microsoft PowerPoint or LibreOffice Impress.")
